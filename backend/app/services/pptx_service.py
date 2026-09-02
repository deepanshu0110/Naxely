import io
import logging
import os
import tempfile
from datetime import datetime

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from app.services.pdf_service import _trend_qualifier

logger = logging.getLogger(__name__)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)
CONTENT_W = SLIDE_W - 2 * MARGIN
TITLE_H = Inches(0.7)

# Ledger-aligned palette — matches pdf_service.py tokens (INK #1A1D24, PAPER_BG #FAFAF7, DELTA_UP #5B7C55, RUST #A8481F)
COLOR_INK = RGBColor(0x1A, 0x1D, 0x24)  # aligned to PDF INK (was 0x14,0x13,0x1F)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_PAPER_BG = RGBColor(0xFA, 0xFA, 0xF7)
COLOR_MINT = RGBColor(0x0E, 0x9F, 0x6E)  # legacy, kept for sentiment/priority fallback
COLOR_RED = RGBColor(0xEF, 0x44, 0x44)  # legacy
COLOR_AMBER = RGBColor(0xD9, 0x7A, 0x34)
COLOR_GRAY = RGBColor(0x6B, 0x72, 0x80)
COLOR_RUST = RGBColor(0xA8, 0x48, 0x1F)  # PDF RUST
COLOR_DELTA_UP = RGBColor(0x5B, 0x7C, 0x55)  # PDF DELTA_UP
COLOR_B45309 = RGBColor(0xB4, 0x53, 0x09)  # PDF medium amber (distinct from COLOR_AMBER brand #D97A34)

SENTIMENT_COLOR = {
    "positive": COLOR_MINT,
    "negative": COLOR_RED,
    "neutral": COLOR_AMBER,
}

PRIORITY_COLOR = {
    "high": COLOR_RUST,      # aligned to PDF #A8481F (was COLOR_RED #EF4444)
    "medium": COLOR_B45309,  # aligned to PDF #B45309 (was COLOR_AMBER #D97A34)
    "low": COLOR_DELTA_UP,   # aligned to PDF #5B7C55 (was COLOR_MINT #0E9F6E)
}

TREND_ARROW = {"up": "\u2191", "increasing": "\u2191", "down": "\u2193", "decreasing": "\u2193", "neutral": "\u2192", "flat": "\u2192", "stable": "\u2192"}


def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    if len(h) != 6:
        return COLOR_MINT
    return RGBColor.from_string(h)


def _blend_toward_white(hex_color: str, alpha: float) -> RGBColor:
    """Solid pre-blend of `hex_color` toward white — opaque RGB equivalent of
    pdf_service._brand_tint(hex, alpha) which returns Color with alpha. Uses
    r = c*alpha + 255*(1-alpha) so pptx shape fills need no transparency."""
    h = hex_color.lstrip("#")
    if len(h) != 6:
        return COLOR_MINT
    try:
        r = int(h[0:2], 16)
        g = int(h[2:4], 16)
        b = int(h[4:6], 16)
    except ValueError:
        return COLOR_MINT
    br = int(round(r * alpha + 255 * (1 - alpha)))
    bg = int(round(g * alpha + 255 * (1 - alpha)))
    bb = int(round(b * alpha + 255 * (1 - alpha)))
    br = max(0, min(255, br))
    bg = max(0, min(255, bg))
    bb = max(0, min(255, bb))
    return RGBColor(br, bg, bb)


def _blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def _add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, italic=False, color=COLOR_INK, font_name="IBM Plex Sans", word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name


def _fill_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()


def _slide_title(slide, text, brand_color):
    _add_textbox(slide, MARGIN, MARGIN, CONTENT_W, TITLE_H, text, font_size=28, bold=True, color=COLOR_INK, font_name="Fraunces")
    _add_rect(slide, MARGIN, MARGIN + TITLE_H, CONTENT_W, Inches(0.04), fill_color=brand_color)


def _build_cover(prs, config, user_data, brand_color):
    slide = _blank_slide(prs)
    # Light ledger background — matches PDF PAPER_BG #FAFAF7
    _fill_background(slide, COLOR_PAPER_BG)

    # Tinted hero panel — solid pre-blend of brand_color toward white at 0.12 (same alpha PDF _CoverHeroStat uses)
    brand_hex = user_data.get("brand_color") or "#0E9F6E"
    # derive hex from RGBColor if user_data hex missing or mismatched
    try:
        derived_hex = f"{brand_color[0]:02X}{brand_color[1]:02X}{brand_color[2]:02X}"
        if derived_hex.lower() != brand_hex.lstrip("#").lower():
            brand_hex = f"#{derived_hex}"
    except Exception:
        pass
    panel_fill = _blend_toward_white(brand_hex, 0.12)
    _add_rect(slide, MARGIN, Inches(1.9), CONTENT_W, Inches(3.4), fill_color=panel_fill)

    # Motif echo — 3 ascending rectangles in brand_color (python-pptx shapes, not ReportLab Flowables)
    motif_w = Inches(0.18)
    motif_gap = Inches(0.05)
    motif_heights = [Inches(0.26), Inches(0.38), Inches(0.52)]
    total_motif_w = 3 * motif_w + 2 * motif_gap
    motif_start_x = SLIDE_W / 2 - total_motif_w / 2
    motif_baseline = Inches(1.55) + max(motif_heights)
    for i, h in enumerate(motif_heights):
        x = motif_start_x + i * (motif_w + motif_gap)
        y = motif_baseline - h
        _add_rect(slide, x, y, motif_w, h, fill_color=brand_color)

    # Keep existing logo-fetch-and-place logic exactly as-is
    logo_url = user_data.get("logo_url")
    if logo_url:
        try:
            import httpx
            resp = httpx.get(logo_url, timeout=5)
            if resp.status_code == 200:
                tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
                tmp.write(resp.content)
                tmp.close()
                slide.shapes.add_picture(tmp.name, MARGIN, MARGIN, width=Inches(1.5))
                os.unlink(tmp.name)
        except Exception as e:
            logger.warning(f"pptx_service: logo fetch failed: {e}")

    company = user_data.get("company_name") or "Report"
    _add_textbox(slide, MARGIN, Inches(2.2), CONTENT_W, Inches(1.0), company, font_size=36, bold=True, color=COLOR_INK, font_name="Fraunces")

    title = config.get("title") or "Naxely Report"
    _add_textbox(slide, MARGIN, Inches(3.3), CONTENT_W, Inches(0.8), title, font_size=24, color=COLOR_INK)

    date_str = datetime.now().strftime("%B %d, %Y")
    _add_textbox(slide, MARGIN, Inches(4.2), CONTENT_W, Inches(0.5), date_str, font_size=16, color=COLOR_INK)

    _add_textbox(slide, SLIDE_W - Inches(3), SLIDE_H - Inches(0.5), Inches(2.4), Inches(0.4), "Generated by Naxely \u00b7 naxely.com", font_size=10, color=COLOR_GRAY)


def _build_executive_summary(prs, summary, brand_color):
    slide = _blank_slide(prs)
    _slide_title(slide, "Executive Summary", brand_color)
    _add_textbox(slide, MARGIN, MARGIN + TITLE_H + Inches(0.2), CONTENT_W, Inches(5.5), summary, font_size=18, color=COLOR_INK, word_wrap=True)


def _build_kpi_slide(prs, kpis, brand_color):
    CARD_W = Inches(5.8)
    CARD_H = Inches(1.6)
    COL_GAP = Inches(0.5)
    ROW_GAP = Inches(0.2)
    ACCENT_W = Inches(0.08)

    # Resolve brand_hex for the blended card fill (opaque equivalent of pdf _KPICard alpha 0.07)
    try:
        brand_hex = f"#{brand_color[0]:02X}{brand_color[1]:02X}{brand_color[2]:02X}"
    except Exception:
        brand_hex = "#0E9F6E"
    card_fill = _blend_toward_white(brand_hex, 0.07)  # matches pdf_service _KPICard tint alpha 0.07

    for page_start in range(0, len(kpis), 6):
        slide = _blank_slide(prs)
        _slide_title(slide, "Key Performance Indicators", brand_color)

        page_kpis = kpis[page_start:page_start + 6]
        for i, kpi in enumerate(page_kpis):
            col = i % 2
            row = i // 2
            left = MARGIN + col * (CARD_W + COL_GAP)
            top = MARGIN + TITLE_H + Inches(0.3) + row * (CARD_H + ROW_GAP)

            # Trend color derived from trend_pct sign — matches pdf_service _KPICard / _KPIRow
            # (DELTA_UP #5B7C55 for >=0, RUST #A8481F for <0). Falls back to trend string if trend_pct missing.
            trend_pct = kpi.get("trend_pct")
            if trend_pct is not None:
                try:
                    pct = float(trend_pct)
                except Exception:
                    pct = 0
                if pct >= 0:
                    card_color = COLOR_DELTA_UP
                    arrow = "\u2191"
                else:
                    card_color = COLOR_RUST
                    arrow = "\u2193"
            else:
                trend = kpi.get("trend", "neutral")
                if trend in ("up", "increasing"):
                    card_color = COLOR_DELTA_UP
                    arrow = "\u2191"
                elif trend in ("down", "decreasing"):
                    card_color = COLOR_RUST
                    arrow = "\u2193"
                else:
                    card_color = COLOR_AMBER
                    arrow = "\u2192"

            _add_rect(slide, left, top, CARD_W, CARD_H, fill_color=card_fill)
            _add_rect(slide, left, top, ACCENT_W, CARD_H, fill_color=card_color)

            _add_textbox(slide, left + ACCENT_W + Inches(0.15), top + Inches(0.1), CARD_W - ACCENT_W - Inches(0.2), Inches(0.4), kpi.get("name", ""), font_size=14, bold=True, color=COLOR_INK)

            qualifier = _trend_qualifier(kpi.get('trend_label', ''))
            if trend_pct is not None:
                # Keep legacy display format: pct:.1f% (no leading + for positive, minus shown for negative) to match existing test expectations
                trend_text = f"{pct:.1f}%" + (f" ({qualifier})" if qualifier else '')
            else:
                # fallback: show raw trend_pct 0 if missing
                trend_text = f"{kpi.get('trend_pct', 0):.1f}%" + (f" ({qualifier})" if qualifier else '')
                arrow = TREND_ARROW.get(kpi.get("trend", "neutral"), arrow)
            value_text = f"{kpi.get('value', '')}  {arrow} {trend_text}"
            _add_textbox(slide, left + ACCENT_W + Inches(0.15), top + Inches(0.55), CARD_W - ACCENT_W - Inches(0.2), Inches(0.6), value_text, font_size=24, bold=True, color=card_color, font_name="IBM Plex Mono")


def _build_chart_slides(prs, chart_paths, brand_color):
    for item in chart_paths:
        if isinstance(item, tuple):
            path = item[0] if len(item) > 0 else ""
            metric = item[1] if len(item) > 1 else "Chart"
            caption = item[2] if len(item) > 2 else None
            insight_title = item[3] if len(item) > 3 else None
        else:
            path, metric, caption, insight_title = item, "Chart", None, None

        if not os.path.exists(path):
            logger.warning(f"pptx_service: chart not found, skipping: {path}")
            continue

        # Prefer insight-driven title when available, fallback to raw metric
        display_title = insight_title.strip() if isinstance(insight_title, str) and insight_title.strip() else metric
        slide = _blank_slide(prs)
        _slide_title(slide, display_title, brand_color)

        try:
            pic_top = MARGIN + TITLE_H + Inches(0.2)
            # If caption exists, reserve space at bottom so it doesn't overlap chart
            caption_text = None
            if isinstance(caption, str) and caption.strip():
                # Strip <b> markup from pdf caption; keep plain text for pptx
                import re as _re
                caption_text = _re.sub(r'<[^>]+>', '', caption).strip()
                if caption_text == "":
                    caption_text = None
            if caption_text:
                caption_h = Inches(0.35)
                pic_h = SLIDE_H - pic_top - caption_h - Inches(0.15) - MARGIN
                slide.shapes.add_picture(path, MARGIN, pic_top, width=CONTENT_W, height=pic_h)
                # Caption as smaller muted text below chart
                caption_top = pic_top + pic_h + Inches(0.12)
                _add_textbox(slide, MARGIN, caption_top, CONTENT_W, caption_h, caption_text, font_size=10, italic=True, color=COLOR_GRAY, font_name="IBM Plex Sans", word_wrap=True)
            else:
                pic_h = SLIDE_H - pic_top - MARGIN
                slide.shapes.add_picture(path, MARGIN, pic_top, width=CONTENT_W, height=pic_h)
        except Exception as e:
            logger.warning(f"pptx_service: failed to embed chart {path}: {e}")


def _build_insights_slides(prs, insights, brand_color):
    CARD_W = CONTENT_W
    CARD_H = Inches(1.55)
    ACCENT_W = Inches(0.07)
    GAP = Inches(0.15)

    # PDF-aligned priority hex — matches _InsightCard in pdf_service.py
    PRIORITY_HEX = {
        "high": "#A8481F",
        "medium": "#B45309",
        "low": "#5B7C55",
    }

    for page_start in range(0, len(insights), 3):
        slide = _blank_slide(prs)
        _slide_title(slide, "AI Insights", brand_color)

        page_insights = insights[page_start:page_start + 3]
        for i, ins in enumerate(page_insights):
            top = MARGIN + TITLE_H + Inches(0.25) + i * (CARD_H + GAP)
            priority = str(ins.get("priority", "medium")).lower().strip() or "medium"
            # PDF fallback is medium (#B45309) for unknown priority
            accent_color = PRIORITY_COLOR.get(priority, COLOR_B45309)
            priority_hex = PRIORITY_HEX.get(priority, "#B45309")
            # Reuse _blend_toward_white helper from Part 1 — opaque equivalent of pdf _brand_tint at 0.12 (pdf _InsightCard alpha)
            card_fill = _blend_toward_white(priority_hex, 0.12)

            _add_rect(slide, MARGIN, top, CARD_W, CARD_H, fill_color=card_fill)
            _add_rect(slide, MARGIN, top, ACCENT_W, CARD_H, fill_color=accent_color)

            x = MARGIN + ACCENT_W + Inches(0.15)
            w = CARD_W - ACCENT_W - Inches(0.2)

            header = f"{ins.get('kpi', '')}  [{priority.upper()}]"
            _add_textbox(slide, x, top + Inches(0.1), w, Inches(0.30), header, font_size=13, bold=True, color=COLOR_INK)
            # number headline — present in ai_service.generate_nra_insights and pdf _InsightCard, but previously missing in pptx
            number = str(ins.get("number", "") or "").strip()
            if number:
                _add_textbox(slide, x, top + Inches(0.38), w, Inches(0.28), number, font_size=11, bold=True, color=COLOR_INK, font_name="IBM Plex Mono")
            _add_textbox(slide, x, top + Inches(0.68), w, Inches(0.38), ins.get("reason", ""), font_size=11, color=COLOR_GRAY)
            _add_textbox(slide, x, top + Inches(1.08), w, Inches(0.35), f"\u2192 {ins.get('action', '')}", font_size=11, color=COLOR_DELTA_UP, italic=True)


def _build_data_table_slides(prs, df, brand_color, config):
    """Paginated Data Table — native add_table grid, zebra striping, colored Status text.
    Rows per slide: 12 data rows (header + 12 = 13 rows total, ~4in tall within 5.4in usable).
    Total cap: 30 rows (3 slides max) — deck-appropriate vs PDF's 50-row print cap."""
    if df is None or df.empty:
        return
    display_df = df.copy()
    # Deck-appropriate total cap: 30 rows → at most 3 slides (12+12+6); PDF's 50 would span 5 slides, too dense for talk-over medium
    TOTAL_CAP = 30
    if len(display_df) > TOTAL_CAP:
        display_df = display_df.head(TOTAL_CAP)
    ROWS_PER_SLIDE = 12  # data rows per slide; header is extra

    # Identify column types for formatting / alignment (mirrors pdf_service)
    numeric_cols = set(c for c in display_df.columns if pd.api.types.is_numeric_dtype(display_df[c]))
    date_cols = set(c for c in display_df.columns if pd.api.types.is_datetime64_any_dtype(display_df[c]))

    # Status column detection — case-insensitive, free-text values (not enum)
    status_col = None
    for c in display_df.columns:
        if str(c).lower() == "status":
            status_col = c
            break

    MAX_CELL_CHARS = 18

    def _truncate_words(text: str, limit: int) -> str:
        if len(text) <= limit:
            return text
        words = text.split(" ")
        for n in range(len(words), 0, -1):
            cand = " ".join(words[:n])
            if len(cand) <= limit - 1:
                return cand + "\u2026"
        return text[: limit - 1] + "\u2026"

    def _fmt(val, col):
        if pd.isna(val):
            return ""
        if col in numeric_cols:
            if isinstance(val, float):
                return f"{val:,.2f}" if val != int(val) else f"{int(val):,}"
            return f"{val:,}"
        text = str(val)
        return _truncate_words(text, MAX_CELL_CHARS)

    headers = [str(c).upper() for c in display_df.columns]

    # Column widths — proportional to header + sample content length (approx, not pixel-identical to PDF's pdfmetrics logic, but legible)
    MIN_W = Inches(0.9)
    MAX_W = Inches(1.9)
    col_widths = []
    for col in display_df.columns:
        if str(col).lower() == "status":
            # Fixed width for Status to keep colored text legible; 68-80pt PDF range ≈ 0.94-1.11in
            col_widths.append(Inches(1.05))
        else:
            samples = [_fmt(v, col) for v in display_df[col].dropna().head(20)]
            max_len = max([len(str(col).upper())] + [len(s) for s in samples] or [0])
            # ~0.11in per char at 9pt mono + 0.30in padding; clamp to MIN/MAX
            w = Inches(max_len * 0.11 + 0.30)
            w = max(MIN_W, min(MAX_W, w))
            col_widths.append(w)
    total = sum(col_widths)
    if total > CONTENT_W:
        scale = CONTENT_W / total
        col_widths = [int(w * scale) for w in col_widths]
    else:
        col_widths = [int(w) for w in col_widths]

    # Pagination — mirror KPI/Insight for page_start in range(0, len(df), 6)
    for page_start in range(0, len(display_df), ROWS_PER_SLIDE):
        chunk = display_df.iloc[page_start : page_start + ROWS_PER_SLIDE]
        n_data = len(chunk)
        rows = n_data + 1  # header + data
        cols = len(display_df.columns)

        slide = _blank_slide(prs)
        # Section title with pagination suffix when multiple slides
        title = "Data Table" if len(display_df) <= ROWS_PER_SLIDE else f"Data Table  ({page_start + 1}–{page_start + n_data} of {len(display_df)})"
        _slide_title(slide, title, brand_color)

        left = MARGIN
        top = MARGIN + TITLE_H + Inches(0.2)
        # Use CONTENT_W for table width; height will be set via row heights
        table_height = Inches(0.35) + n_data * Inches(0.28)
        # Clamp to usable vertical (5.4in); if overflow, row height already accounts
        table_shape = slide.shapes.add_table(rows, cols, left, top, CONTENT_W, table_height)
        table = table_shape.table

        # Apply column widths
        for j, w in enumerate(col_widths):
            table.columns[j].width = w
        # Row heights — header slightly taller
        table.rows[0].height = Inches(0.35)
        for r in range(1, rows):
            table.rows[r].height = Inches(0.28)

        # Colors — zebra matching PDF spirit: even #F2F1EB, odd #FAFAF7; header muted on light
        even_fill = RGBColor(0xF2, 0xF1, 0xEB)
        odd_fill = RGBColor(0xFA, 0xFA, 0xF7)
        header_fill = RGBColor(0xFF, 0xFF, 0xFF)  # white header for contrast
        header_text_color = COLOR_GRAY  # muted, matches PDF header

        # Header row
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = h
            # Cell fill
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill
            # Font
            para = cell.text_frame.paragraphs[0]
            para.alignment = 2  # CENTER (PP_PARAGRAPH_ALIGNMENT.CENTER = 2)
            run = para.runs[0]
            run.font.size = Pt(9)
            run.font.bold = True
            run.font.color.rgb = header_text_color
            run.font.name = "IBM Plex Sans"
            cell.vertical_anchor = 3  # MIDDLE
            cell.text_frame.word_wrap = True
            # Header bottom border via line? python-pptx tables have limited border API; rely on fill contrast + slide title rule

        # Data rows
        for i, (_, row) in enumerate(chunk.iterrows()):
            r = i + 1  # table row index (0 is header)
            # Zebra
            fill_rgb = even_fill if r % 2 == 0 else odd_fill
            for j, col in enumerate(display_df.columns):
                cell = table.cell(r, j)
                val = row[col]
                text = _fmt(val, col)
                cell.text = text
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill_rgb
                para = cell.text_frame.paragraphs[0]
                # Alignment per column type
                if status_col is not None and str(col).lower() == "status":
                    para.alignment = 2  # CENTER
                elif col in numeric_cols:
                    para.alignment = 4  # RIGHT
                elif col in date_cols:
                    para.alignment = 2  # CENTER
                else:
                    para.alignment = 1  # LEFT
                run = para.runs[0]
                run.font.size = Pt(9)
                run.font.name = "IBM Plex Mono" if col in numeric_cols or col in date_cols else "IBM Plex Sans"
                run.font.bold = False
                # Status colored-text treatment — case-insensitive value, fallback to neutral gray
                if status_col is not None and str(col).lower() == "status":
                    val_str = str(val) if val is not None else ""
                    if val_str.lower().strip() == "billed":
                        run.font.color.rgb = COLOR_DELTA_UP
                        run.font.bold = True
                    else:
                        run.font.color.rgb = COLOR_GRAY
                        # keep bold False for pending/other
                    run.font.size = Pt(8.5)
                    run.font.name = "IBM Plex Sans"
                else:
                    run.font.color.rgb = COLOR_INK
                cell.vertical_anchor = 3  # MIDDLE
                cell.text_frame.word_wrap = True
                # Ensure margin inside cell
                cell.margin_left = Inches(0.06)
                cell.margin_right = Inches(0.06)
                cell.margin_top = Inches(0.02)
                cell.margin_bottom = Inches(0.02)


def _build_recommendations_slides(prs, recommendations, brand_color):
    """Stacked recommendation cards — tinted panel 0.12, left accent 5pt, numbered badge 01-05.
    Height estimated per text length (chars-per-line at chosen font width), pagination when overflow."""
    if not recommendations:
        return
    try:
        brand_hex = f"#{brand_color[0]:02X}{brand_color[1]:02X}{brand_color[2]:02X}"
    except Exception:
        brand_hex = "#0E9F6E"
    tint = _blend_toward_white(brand_hex, 0.12)  # matches pdf _RecommendationCard alpha 0.12, also cover/insights
    ACCENT_W = Inches(0.07)  # 5pt left strip, same as insights
    BADGE_SIZE = Inches(0.42)  # 30pt circle
    BADGE_X = Inches(0.12)
    TEXT_GAP = Inches(0.15)
    CARD_GAP = Inches(0.12)
    # Text geometry inside card
    text_x_offset = BADGE_X + BADGE_SIZE + TEXT_GAP  # from card left
    text_width = CONTENT_W - text_x_offset - Inches(0.12)  # right padding
    # Estimation: IBM Plex Sans 11pt ~6.6pt per char average, word-wrap
    FONT_SIZE = 11
    LINE_HEIGHT = Pt(14)  # PDF leading 14pt
    AVG_CHAR_PT = FONT_SIZE * 0.6
    text_width_pt = text_width / 12700  # EMU to pt: 12700 EMU = 1pt (914400/72)
    # Fallback if calc off: CONTENT_W ~12.13in => text_width ~11.3in => 813pt => ~123 chars/line at 6.6pt/char
    chars_per_line = max(30, int(text_width_pt / AVG_CHAR_PT)) if text_width_pt else 80

    def _estimate_lines(text: str) -> int:
        if not text:
            return 1
        # Rough word-wrap: split by words, greedy fill per line
        words = text.split()
        lines = 1
        cur_len = 0
        for w in words:
            # +1 for space except first word per line
            need = len(w) + (1 if cur_len else 0)
            if cur_len + need <= chars_per_line:
                cur_len += need
            else:
                lines += 1
                cur_len = len(w)
        return max(1, lines)

    def _estimate_card_height(text: str) -> int:
        lines = _estimate_lines(text)
        text_h = lines * LINE_HEIGHT  # EMU
        # Card padding top+bottom 12pt each (PDF 24) + ensure at least badge+margin
        card_h = max(Inches(0.70), text_h + Inches(0.24))
        # Also ensure badge fits (BADGE_SIZE + 0.2in margin)
        card_h = max(card_h, BADGE_SIZE + Inches(0.20))
        return int(card_h)

    # Pagination state
    slide = None
    cur_top = None
    usable_top = MARGIN + TITLE_H + Inches(0.25)
    usable_bottom = SLIDE_H - MARGIN
    usable_height = usable_bottom - usable_top

    for idx, rec in enumerate(recommendations, 1):
        text = str(rec).strip()
        if not text:
            continue
        card_h = _estimate_card_height(text)
        # Start new slide if needed or no slide yet
        if slide is None or (cur_top + card_h > usable_bottom):
            slide = _blank_slide(prs)
            # Title with pagination suffix when multiple slides
            ttl = "Recommendations" if len(recommendations) <= 5 else f"Recommendations ({idx}–{min(idx+4, len(recommendations))} of {len(recommendations)})"
            # Simpler: show total count only when paginated; reuse single title style
            if len(recommendations) > 4 and idx > 1:
                # For overflow case, add slide-number hint in title? Keep plain
                pass
            _slide_title(slide, "Recommendations", brand_color)
            cur_top = usable_top

        # Card background tinted panel + left accent
        card_left = MARGIN
        card_width = CONTENT_W
        _add_rect(slide, card_left, cur_top, card_width, card_h, fill_color=tint)
        _add_rect(slide, card_left, cur_top, ACCENT_W, card_h, fill_color=brand_color)

        # Numbered badge — oval (circle) brand_color fill
        badge_left = card_left + BADGE_X
        badge_top = cur_top + (card_h - BADGE_SIZE) / 2
        # oval shape type 9
        badge_shape = slide.shapes.add_shape(9, int(badge_left), int(badge_top), int(BADGE_SIZE), int(BADGE_SIZE))
        badge_shape.fill.solid()
        badge_shape.fill.fore_color.rgb = brand_color
        badge_shape.line.fill.background()
        # Badge number centered white bold
        num_text = str(idx).zfill(2)
        # Textbox covering badge for centered number
        txBox = slide.shapes.add_textbox(int(badge_left), int(badge_top), int(BADGE_SIZE), int(BADGE_SIZE))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = 2  # CENTER
        p.text = num_text
        run = p.runs[0]
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = COLOR_WHITE
        run.font.name = "Fraunces"
        # Vertical centering via paragraph space — pptx vertical_anchor middle
        tf.vertical_anchor = 3  # MIDDLE

        # Recommendation text
        text_left = card_left + text_x_offset
        # Text box height slightly smaller than card to allow padding
        text_top = cur_top + Inches(0.12)
        text_h = card_h - Inches(0.24)
        _add_textbox(slide, int(text_left), int(text_top), int(text_width), int(text_h), text, font_size=FONT_SIZE, bold=False, color=COLOR_INK, font_name="IBM Plex Sans", word_wrap=True)

        cur_top += card_h + CARD_GAP


def _build_anomaly_slide(prs, anomalies, brand_color):
    slide = _blank_slide(prs)
    _slide_title(slide, "Anomaly Flags", brand_color)

    bullet_top = MARGIN + TITLE_H + Inches(0.3)
    for i, anomaly in enumerate(anomalies):
        msg = anomaly.get("message", "")
        _add_textbox(slide, MARGIN + Inches(0.2), bullet_top + i * Inches(0.55), CONTENT_W - Inches(0.2), Inches(0.5), f"\u2022 {msg}", font_size=14, color=COLOR_INK, word_wrap=True)


def _build_footer_slide(prs, brand_color):
    slide = _blank_slide(prs)
    _fill_background(slide, brand_color)
    _add_textbox(slide, MARGIN, Inches(3.2), CONTENT_W, Inches(1.0), "Generated by Naxely", font_size=32, bold=True, color=COLOR_WHITE, font_name="Fraunces")
    _add_textbox(slide, MARGIN, Inches(4.3), CONTENT_W, Inches(0.5), "naxely.com", font_size=18, color=COLOR_WHITE)


def generate_pptx(
    df: pd.DataFrame,
    chart_paths: list,
    ai_content: dict,
    config: dict,
    user_data: dict,
) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    brand_hex = user_data.get("brand_color") or "#0E9F6E"
    brand_color = _hex_to_rgb(brand_hex)
    ai_skipped = config.get("_ai_skipped", False)
    sections = config.get("sections", [])

    _build_cover(prs, config, user_data, brand_color)

    summary = ai_content.get("summary")
    if summary and not ai_skipped:
        _build_executive_summary(prs, summary, brand_color)
    elif ai_skipped:
        slide = _blank_slide(prs)
        _slide_title(slide, "AI Analysis", brand_color)
        _add_textbox(slide, MARGIN, MARGIN + TITLE_H + Inches(0.3), CONTENT_W, Inches(2.0),
                      "AI analysis was skipped for this report.\nRe-generate with a valid AI provider key to include insights.",
                      font_size=16, color=COLOR_GRAY)

    kpis = config.get("_precomputed_kpis") or []
    if kpis:
        _build_kpi_slide(prs, kpis, brand_color)

    if chart_paths and ("charts" in sections or not sections):
        _build_chart_slides(prs, chart_paths, brand_color)

    insights = ai_content.get("insights") or []
    if insights and not ai_skipped and ("insights" in sections or not sections):
        _build_insights_slides(prs, insights, brand_color)

    anomalies = ai_content.get("anomalies") or []
    if anomalies and ("anomalies" in sections or not sections):
        _build_anomaly_slide(prs, anomalies, brand_color)

    if "data_table" in sections:
        _build_data_table_slides(prs, df, brand_color, config)

    # Recommendations — same gate as pdf_service.py: ('insights' or 'executive_summary') and not ai_skipped
    show_recommendations = (("insights" in sections) or ("executive_summary" in sections)) and not ai_skipped
    if show_recommendations:
        recommendations = ai_content.get("recommendations") or []
        # Only build if list non-empty; _build handles empty early return
        if recommendations:
            _build_recommendations_slides(prs, recommendations, brand_color)

    _build_footer_slide(prs, brand_color)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
