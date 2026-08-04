"""PDF/PPTX KPI parity regression tests.

Regression for the Summit & Holt bug class: the PPTX export route used to
recompute KPI trends on a date-blind code path (raw first-to-last) while the
PDF showed monthly-aggregate trends, so the same report displayed different
trend_pct values in the two exports (e.g. Revenue +101.8% in the PDF vs
+195.1% in the PPTX). The fix persists the PDF's precomputed KPI dict into the
stored report config (report_service._build_stored_config) and the PPTX route
reuses it instead of recomputing.

These tests run the REAL code paths end-to-end:
- PDF side: parse_csv -> _process_csv -> _compute_kpi_data ->
  _build_stored_config (the exact persistence the pipeline performs) ->
  build_sync, producing a real PDF whose text is extracted with PyMuPDF.
- PPTX side: the real export_report_pptx route (storage/DB mocked only),
  producing a real .pptx whose text is extracted with python-pptx.
"""

import io
import json
import os
import re
import sys

sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..'))

import numpy as np
import pandas as pd
import pytest
from unittest.mock import patch, MagicMock, AsyncMock

from app.services.data_service import parse_csv
from app.services.report_service import _process_csv, _build_stored_config
from app.services.pdf_service import _compute_kpi_data

KPI_NAMES = ["Total Units Sold", "Avg Unit Price", "Total Revenue", "Total Profit"]

AI_EMPTY = {"summary": None, "insights": [], "anomalies": [], "trends": []}

REPORT_ID = "parity-report-001"
USER_ID = "agency-user-parity"


# ── Dataset: Summit & Holt shape ──────────────────────────────────────────────
# 48 rows, 12 unique dates (2026-01-01..2026-12-01), 4 rows per month.
# Unit_Price has exactly 2 distinct values {29, 79} alternating per row, so its
# monthly sums are constant (216 == 216) -> genuinely flat month-over-month.
# Units/Revenue/Profit grow month to month, so monthly-aggregate trend_pct
# (97.8 / 101.8 / 132.6) differs sharply from raw first-to-last
# (8.3 / 195.1 / 267.1) - the exact shape that surfaced the original bug.

def _build_demo_csv_bytes():
    months, rpm = 12, 4
    dates, units, rev, profit = [], [], [], []
    for m in range(months):
        dates += [f"2026-{m+1:02d}-01"] * rpm
        if m == 0:
            units += [12, 5, 8, 21]            # monthly sum 46
            rev += [348, 395, 232, 540]        # monthly sum 1515
            profit += [228, 245, 142, 285]     # monthly sum 900
        elif m == 11:
            units += [22, 24, 32, 13]          # monthly sum 91
            rev += [620, 690, 720, 1027]       # monthly sum 3057
            profit += [420, 430, 406, 837]     # monthly sum 2093
        else:
            base = int(round(1515 + (3057 - 1515) * m / 11)) // 4
            units += [10 + m, 12 + m, 9 + m, 13 + m]
            rev += [base] * 4
            profit += [int(round(0.62 * base))] * 4

    rng = np.random.default_rng(7)
    N = months * rpm
    df = pd.DataFrame({
        "Date": dates,
        "Region": rng.choice(["North", "South", "East", "West"], N),
        "Salesperson": rng.choice(["Alice", "Bob", "Carol", "Dave"], N),
        "Product": rng.choice(["Pro Plan", "Agency Plan"], N),
        "Category": ["SaaS"] * N,
        "Units_Sold": units,
        "Unit_Price": [29, 79, 29, 79] * months,
        "Revenue": rev,
        "Cost": [r - p for r, p in zip(rev, profit)],
        "Profit": profit,
        "Customer_Satisfaction": np.round(rng.uniform(3.8, 4.9, N), 1),
        "Leads_Generated": rng.integers(20, 60, N),
        "Conversion_Rate": np.round(rng.uniform(0.1, 0.4, N), 2),
    })
    return df.to_csv(index=False).encode("utf-8")


# ── PDF side (mirrors run_report_pipeline data -> kpi -> persist) ─────────────

def _pdf_side(csv_bytes):
    """Run the pipeline's real data/KPI/persistence steps and build a real PDF.
    Returns (pdf_text, kpis, stored_config) where stored_config is exactly what
    the pipeline now persists into the reports row."""
    from app.services.pdf_service import build_sync

    df = parse_csv(csv_bytes)
    config = {
        "upload_id": "upl-parity-001",
        "title": "Summit & Holt — Monthly Performance",
        "sections": ["kpi_overview", "data_table"],
        "template_type": "professional",
    }
    df, df_norm = _process_csv(df, config)
    assert config["date_column"] == "Date"

    metric_cols = config.get("metric_columns") or [
        c for c in df_norm.columns if pd.api.types.is_numeric_dtype(df_norm[c])
    ]
    config["metric_columns"] = metric_cols[:5]

    kpis = _compute_kpi_data(df_norm, config, AI_EMPTY, "#6366F1")
    stored_config = _build_stored_config(config, kpis)

    pdf_config = dict(config)
    pdf_config["report_id"] = REPORT_ID
    pdf_config["_precomputed_kpis"] = kpis
    pdf_config["_ai_skipped"] = False
    user_data = {"brand_color": "#6366F1", "tier": "agency", "logo_url": None, "company_name": "Acme Corp"}

    pdf_path = build_sync(df, [], AI_EMPTY, pdf_config, user_data)
    import fitz
    try:
        doc = fitz.open(pdf_path)
        text = "".join(page.get_text() for page in doc)
        doc.close()
    finally:
        try:
            os.unlink(pdf_path)
        except OSError:
            pass
    return text, kpis, stored_config


# ── PPTX side (real export_report_pptx route; DB/storage mocked) ──────────────

async def _pptx_side(csv_bytes, stored_config):
    """Call the real export route and return the produced .pptx bytes."""
    from app.api.routes.reports import export_report_pptx
    from app.models.user import User

    user = User()
    user.id = USER_ID
    user.tier = "agency"
    user.email = "agency@parity.test"

    report_row = {
        "id": REPORT_ID,
        "user_id": USER_ID,
        "title": "Summit & Holt — Monthly Performance",
        "template_type": "professional",
        "status": "completed",
        "source_type": "csv",
        "source_filename": "naxely_demo.csv",
        "config": json.dumps(stored_config),
        "ai_summary": None,
        "ai_insights": [],
        "ai_anomalies": [],
        "ai_skipped": False,
    }

    mock_db = MagicMock()

    async def fake_execute(stmt, *a, **kw):
        sql = getattr(stmt, "text", str(stmt))
        result = MagicMock()
        if "FROM users" in sql:
            result.mappings.return_value.first.return_value = {
                "brand_color": None,
                "company_name": "Acme Corp",
                "logo_url": None,
                "tier": "agency",
            }
        else:
            result.mappings.return_value.first.return_value = dict(report_row)
        return result

    mock_db.execute = AsyncMock(side_effect=fake_execute)
    mock_db.commit = AsyncMock()

    storage = MagicMock()
    uploads_bucket = MagicMock()
    uploads_bucket.download = lambda path: csv_bytes
    reports_bucket = MagicMock()
    reports_bucket.upload = lambda *a, **kw: None
    storage.storage.from_.side_effect = (
        lambda bucket: uploads_bucket if bucket == "uploads" else reports_bucket
    )

    with (
        patch("app.api.routes.reports._get_supabase", return_value=storage),
        patch("app.services.report_service.get_upload", new_callable=AsyncMock) as mock_get_upload,
    ):
        mock_get_upload.return_value = {
            "id": "upl-parity-001",
            "filename": "naxely_demo.csv",
            "source_type": "csv",
            "file_url": "naxely_demo.csv",
        }
        response = await export_report_pptx(report_id=REPORT_ID, current_user=user, db=mock_db)

    return response.body


# ── Text extraction helpers ───────────────────────────────────────────────────

def _extract_pdf_kpi_pcts(text):
    """Map KPI name -> trend_pct (float) from PDF card text. KPI cards render
    '↑ +101.8%' (arrow anchored), so the cover hero and summary strip (which
    render '+101.8%' without an arrow) never match."""
    name_re = re.compile("|".join(re.escape(n) for n in KPI_NAMES))
    pct_re = re.compile(r"[↑↓]\s*([+-])\s*(\d+\.\d)\s*%")
    name_pos = [(m.start(), m.group(0)) for m in name_re.finditer(text)]
    out = {}
    for i, (pos, name) in enumerate(name_pos):
        end = name_pos[i + 1][0] if i + 1 < len(name_pos) else len(text)
        m = pct_re.search(text[pos:end])
        if m:
            sign = "-" if m.group(1) == "-" else ""
            out.setdefault(name, float(sign + m.group(2)))
    return out


def _extract_pptx_kpi_pcts(prs):
    """Map KPI name -> trend_pct (float) from PPTX slide text. Each card is a
    name textbox immediately followed by a 'value  ↑ 101.8% (qualifier)'
    textbox, so pair each name with the next arrow pct, then reset."""
    name_re = re.compile("|".join(re.escape(n) for n in KPI_NAMES))
    pct_re = re.compile(r"[↑↓→]\s*([+-]?\d+\.\d)\s*%")
    out = {}
    for slide in prs.slides:
        current = None
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            t = shape.text
            m = name_re.search(t)
            if m:
                current = m.group(0)
            mp = pct_re.search(t)
            if mp and current:
                out[current] = float(mp.group(1))
                current = None
    return out


# ── Tests ─────────────────────────────────────────────────────────────────────

class TestPptxPdfKpiParity:

    @pytest.mark.asyncio
    async def test_pdf_and_pptx_show_same_kpi_trend_pct(self):
        """End-to-end: PDF and PPTX generated from the same report show the
        SAME trend_pct for every metric (monthly-aggregate basis), and the
        route reused the persisted dict rather than recomputing."""
        csv_bytes = _build_demo_csv_bytes()
        pdf_text, kpis, stored_config = _pdf_side(csv_bytes)
        pdf_pcts = _extract_pdf_kpi_pcts(pdf_text)

        pptx_bytes = await _pptx_side(csv_bytes, stored_config)
        from pptx import Presentation
        prs = Presentation(io.BytesIO(pptx_bytes))
        pptx_pcts = _extract_pptx_kpi_pcts(prs)

        assert set(pdf_pcts) == set(KPI_NAMES), f"PDF missing KPI cards: {pdf_pcts}"
        assert set(pptx_pcts) == set(KPI_NAMES), f"PPTX missing KPI cards: {pptx_pcts}"
        for name in KPI_NAMES:
            assert pdf_pcts[name] == pptx_pcts[name], (
                f"{name} diverges: PDF {pdf_pcts[name]} vs PPTX {pptx_pcts[name]}"
            )

        # The shared value is the monthly-aggregate one, not raw first-to-last:
        assert pdf_pcts["Total Units Sold"] == pytest.approx(97.8, abs=0.05)
        assert pdf_pcts["Total Revenue"] == pytest.approx(101.8, abs=0.05)
        assert pdf_pcts["Total Profit"] == pytest.approx(132.6, abs=0.05)

        # Fallback proof: a from-scratch recompute on the stored config (no
        # date_column, no _precomputed_kpis) still yields the OLD divergent raw
        # values — so the route's output matching the PDF proves it reused the
        # persisted dict, and the backward-compat fallback remains functional.
        fallback_config = {k: v for k, v in stored_config.items() if k != "_precomputed_kpis"}
        df = pd.read_csv(io.BytesIO(csv_bytes))
        df_norm = df.copy()
        for col in df_norm.select_dtypes(include=["object"]).columns:
            df_norm[col] = pd.to_numeric(df_norm[col], errors="ignore")
        fallback_kpis = _compute_kpi_data(df_norm, fallback_config, AI_EMPTY, "#6366F1")
        fallback_pcts = {k["name"]: k["trend_pct"] for k in fallback_kpis}
        assert fallback_pcts["Total Revenue"] == pytest.approx(195.1, abs=0.05)
        assert fallback_pcts["Total Profit"] == pytest.approx(267.1, abs=0.05)
        assert fallback_pcts["Avg Unit Price"] == pytest.approx(172.4, abs=0.05)
        assert pptx_pcts["Total Revenue"] == pytest.approx(101.8, abs=0.05)

    @pytest.mark.asyncio
    async def test_flat_two_value_series_is_zero_in_both_exports(self):
        """Regression for the Avg Unit Price case: a column with exactly two
        values {29, 79} spread evenly per month is genuinely flat month-over-
        month — BOTH exports must show 0.0%, never the raw (79-29)/29 = 172.4%
        range."""
        csv_bytes = _build_demo_csv_bytes()
        pdf_text, kpis, stored_config = _pdf_side(csv_bytes)
        pdf_pcts = _extract_pdf_kpi_pcts(pdf_text)

        pptx_bytes = await _pptx_side(csv_bytes, stored_config)
        from pptx import Presentation
        prs = Presentation(io.BytesIO(pptx_bytes))
        pptx_pcts = _extract_pptx_kpi_pcts(prs)

        assert pdf_pcts["Avg Unit Price"] == 0.0
        assert pptx_pcts["Avg Unit Price"] == 0.0
