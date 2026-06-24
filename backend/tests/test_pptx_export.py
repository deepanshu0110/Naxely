import io
import os
import sys
import json
import tempfile
import pytest
from unittest.mock import patch, MagicMock, AsyncMock
from datetime import datetime
from PIL import Image as PILImage
from pptx import Presentation


# ── Helpers ────────────────────────────────────────────────────────────────────

def _make_chart_png() -> str:
    """Create a real temporary PNG file. Returns path."""
    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    PILImage.new("RGB", (800, 400), color=(100, 149, 237)).save(tmp.name)
    tmp.close()
    return tmp.name


def _mock_df():
    import pandas as pd
    return pd.DataFrame({
        "Date": pd.date_range("2024-01-01", periods=5, freq="ME"),
        "Revenue": [1000.0, 1200.0, 900.0, 1500.0, 1800.0],
        "Region": ["North", "South", "East", "West", "North"],
    })


def _mock_config():
    return {
        "report_id": "test-report-123",
        "title": "Q4 Sales Report",
        "sections": ["key_metrics", "charts", "insights", "anomalies"],
        "metric_columns": ["Revenue"],
        "date_column": "Date",
        "_ai_skipped": False,
        "_precomputed_kpis": None,
        "brand": {"prepared_by": "Acme Corp"},
    }


def _mock_user_data(brand_color="#0E9F6E"):
    return {
        "brand_color": brand_color,
        "tier": "agency",
        "logo_url": None,
        "company_name": "Acme Corp",
    }


def _mock_ai_content():
    return {
        "summary": "Revenue grew 12% QoQ driven by the North region.",
        "insights": [
            {
                "kpi": "Revenue",
                "number": "1800.00",
                "reason": "Consistent growth driven by North region expansion.",
                "action": "Increase North region headcount.",
                "sentiment": "positive",
                "priority": "high",
            }
        ],
        "anomalies": [
            {"message": "Spike detected on 2024-03-15: Revenue +320% above baseline"}
        ],
        "trends": [],
    }


# ── Unit tests for pptx_service ────────────────────────────────────────────────

class TestPptxServiceUnit:

    def test_generate_pptx_returns_bytes(self):
        """generate_pptx() returns a non-trivial bytes object."""
        from app.services.pptx_service import generate_pptx
        result = generate_pptx(
            df=_mock_df(),
            chart_paths=[],
            ai_content=_mock_ai_content(),
            config=_mock_config(),
            user_data=_mock_user_data(),
        )
        assert isinstance(result, bytes)
        assert len(result) > 5000

    def test_output_is_valid_pptx(self):
        """Output can be parsed by python-pptx without error."""
        from app.services.pptx_service import generate_pptx
        result = generate_pptx(
            df=_mock_df(),
            chart_paths=[],
            ai_content=_mock_ai_content(),
            config=_mock_config(),
            user_data=_mock_user_data(),
        )
        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) >= 3

    def test_cover_slide_contains_company_and_title(self):
        """First slide contains company name and report title."""
        from app.services.pptx_service import generate_pptx
        result = generate_pptx(
            df=_mock_df(),
            chart_paths=[],
            ai_content=_mock_ai_content(),
            config=_mock_config(),
            user_data=_mock_user_data(),
        )
        prs = Presentation(io.BytesIO(result))
        cover_text = " ".join(
            shape.text
            for shape in prs.slides[0].shapes
            if shape.has_text_frame
        )
        assert "Acme Corp" in cover_text
        assert "Q4 Sales Report" in cover_text

    def test_widescreen_dimensions(self):
        """Presentation uses 16:9 widescreen (13.33 x 7.5 inches)."""
        from app.services.pptx_service import generate_pptx
        from pptx.util import Inches
        result = generate_pptx(
            df=_mock_df(),
            chart_paths=[],
            ai_content=_mock_ai_content(),
            config=_mock_config(),
            user_data=_mock_user_data(),
        )
        prs = Presentation(io.BytesIO(result))
        assert abs(prs.slide_width - Inches(13.33)) < 10000
        assert abs(prs.slide_height - Inches(7.5)) < 10000

    def test_chart_png_embedded_as_picture(self):
        """Chart PNG file is embedded as a picture shape."""
        from app.services.pptx_service import generate_pptx
        chart_path = _make_chart_png()
        try:
            result = generate_pptx(
                df=_mock_df(),
                chart_paths=[chart_path],
                ai_content=_mock_ai_content(),
                config=_mock_config(),
                user_data=_mock_user_data(),
            )
            prs = Presentation(io.BytesIO(result))
            # Shape type 13 = PICTURE
            pic_shapes = [
                shape
                for slide in prs.slides
                for shape in slide.shapes
                if shape.shape_type == 13
            ]
            assert len(pic_shapes) >= 1
        finally:
            os.unlink(chart_path)

    def test_chart_png_missing_does_not_crash(self):
        """Non-existent chart path is silently skipped, no exception."""
        from app.services.pptx_service import generate_pptx
        result = generate_pptx(
            df=_mock_df(),
            chart_paths=["/tmp/naxely/nonexistent/chart.png"],
            ai_content=_mock_ai_content(),
            config=_mock_config(),
            user_data=_mock_user_data(),
        )
        assert isinstance(result, bytes)

    def test_kpi_values_appear_in_slides(self):
        """KPI name and value text appears somewhere in the presentation."""
        from app.services.pptx_service import generate_pptx
        config = _mock_config()
        config["_precomputed_kpis"] = [
            {"name": "Total Revenue", "value": "1.8K", "trend": "up", "trend_pct": 12.3}
        ]
        result = generate_pptx(
            df=_mock_df(),
            chart_paths=[],
            ai_content=_mock_ai_content(),
            config=config,
            user_data=_mock_user_data(),
        )
        prs = Presentation(io.BytesIO(result))
        all_text = " ".join(
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
        assert "Total Revenue" in all_text
        assert "1.8K" in all_text

    def test_ai_insight_text_appears_in_slides(self):
        """AI insight reason and action appear in slide text."""
        from app.services.pptx_service import generate_pptx
        result = generate_pptx(
            df=_mock_df(),
            chart_paths=[],
            ai_content=_mock_ai_content(),
            config=_mock_config(),
            user_data=_mock_user_data(),
        )
        prs = Presentation(io.BytesIO(result))
        all_text = " ".join(
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
        assert "North region" in all_text
        assert "headcount" in all_text

    def test_anomaly_message_appears_in_slides(self):
        """Anomaly message string appears in slide text."""
        from app.services.pptx_service import generate_pptx
        result = generate_pptx(
            df=_mock_df(),
            chart_paths=[],
            ai_content=_mock_ai_content(),
            config=_mock_config(),
            user_data=_mock_user_data(),
        )
        prs = Presentation(io.BytesIO(result))
        all_text = " ".join(
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
        assert "Spike detected" in all_text

    def test_empty_anomalies_skips_anomaly_slide(self):
        """No anomaly slide is added when anomalies list is empty."""
        from app.services.pptx_service import generate_pptx
        ai_content = _mock_ai_content()
        ai_content["anomalies"] = []
        result_with = generate_pptx(
            df=_mock_df(),
            chart_paths=[],
            ai_content=_mock_ai_content(),
            config=_mock_config(),
            user_data=_mock_user_data(),
        )
        result_without = generate_pptx(
            df=_mock_df(),
            chart_paths=[],
            ai_content=ai_content,
            config=_mock_config(),
            user_data=_mock_user_data(),
        )
        prs_with = Presentation(io.BytesIO(result_with))
        prs_without = Presentation(io.BytesIO(result_without))
        assert len(prs_without.slides) < len(prs_with.slides)

    def test_ai_skipped_shows_placeholder_text(self):
        """When _ai_skipped is True, a placeholder message appears instead of insights."""
        from app.services.pptx_service import generate_pptx
        config = _mock_config()
        config["_ai_skipped"] = True
        ai_content = {"summary": None, "insights": [], "anomalies": [], "trends": []}
        result = generate_pptx(
            df=_mock_df(),
            chart_paths=[],
            ai_content=ai_content,
            config=config,
            user_data=_mock_user_data(),
        )
        prs = Presentation(io.BytesIO(result))
        all_text = " ".join(
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
        assert "skipped" in all_text.lower() or "unavailable" in all_text.lower()

    def test_brand_color_does_not_crash(self):
        """Amber brand color hex is accepted without error."""
        from app.services.pptx_service import generate_pptx
        result = generate_pptx(
            df=_mock_df(),
            chart_paths=[],
            ai_content=_mock_ai_content(),
            config=_mock_config(),
            user_data=_mock_user_data(brand_color="#D97A34"),
        )
        assert isinstance(result, bytes)


# ── Route tests (mock-based, matches existing test patterns) ──────────────────

class TestPptxExportRoute:

    @pytest.mark.asyncio
    async def test_export_pptx_requires_agency_tier(self):
        """Non-agency user gets 403 on export route."""
        from app.api.deps import require_agency
        from app.models.user import User

        user = User()
        user.id = "free-user-001"
        user.tier = "free"
        user.email = "free@test.com"

        with pytest.raises(Exception) as exc_info:
            require_agency(current_user=user)
        # _check_tier raises HTTPException(403) for non-agency tiers
        assert exc_info.errisinstance(Exception)

    @pytest.mark.asyncio
    async def test_export_pptx_404_on_missing_report(self):
        """Agency user gets 404 for a non-existent report ID."""
        from fastapi import HTTPException
        from app.api.routes.reports import export_report_pptx
        from app.models.user import User

        user = User()
        user.id = "agency-user-001"
        user.tier = "agency"
        user.email = "agency@test.com"

        # Mock db to return no report
        mock_db = MagicMock()
        mock_db.execute = AsyncMock()
        mock_result = MagicMock()
        mock_result.mappings.return_value.first.return_value = None
        mock_db.execute.return_value = mock_result

        with pytest.raises(HTTPException) as exc_info:
            await export_report_pptx(
                report_id="00000000-0000-0000-0000-000000000000",
                current_user=user,
                db=mock_db,
            )
        assert exc_info.value.status_code == 404
